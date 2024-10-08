from dotenv import load_dotenv
import os 
import re

from flask import Flask, jsonify, request , send_file
from flask_cors import CORS, cross_origin

#LangChain Imports
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyMuPDFLoader
from langchain_community.document_loaders import TextLoader
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain.chains.retrieval import create_retrieval_chain
from langchain.chains.history_aware_retriever import create_history_aware_retriever
from langchain.prompts import PromptTemplate
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.messages import HumanMessage, AIMessage 
from langchain_openai import AzureChatOpenAI
from langchain.tools import StructuredTool
from langchain.chains.conversation.memory import ConversationBufferMemory
from langchain_community.chat_message_histories import ChatMessageHistory

#RAG Imports
from langchain_community.vectorstores.azure_cosmos_db import AzureCosmosDBVectorSearch, CosmosDBSimilarityType, CosmosDBVectorSearchType
from langchain_openai import AzureOpenAIEmbeddings
from pymongo import MongoClient
from urllib.parse import quote_plus


from pptx import Presentation
from pptx.util import Inches, Pt
import json

load_dotenv()
app = Flask(__name__)

# LLM Initialization
llm_api_key = os.getenv('CONST_GPT4o_API_KEY')
llm_api_version = os.getenv('CONST_GPT4o_API_VERSION')
llm_base_url = os.getenv('CONST_GPT4o_ENDPOINT')

llm = AzureChatOpenAI(
                api_key=llm_api_key,
                api_version=llm_api_version,
                base_url=llm_base_url
            )

# Embeddings Initialization
embedding_azure_deployment = os.getenv('CONST_OPENAI_EMBEDDING_DEPLOYMENT')
embedding_api_key = os.getenv('CONST_OPENAI_EMBEDDING_API_KEY')
embedding_model = os.getenv('CONST_OPENAI_EMBEDDING_MODEL_NAME')
embedding_azure_endpoint = os.getenv('CONST_OPENAI_EMBEDDING_ENDPOINT')
embedding_api_version = os.getenv('CONST_OPENAI_EMBEDDING_API_VERSION')

embedding_model= AzureOpenAIEmbeddings(
    azure_deployment=embedding_azure_deployment,
    api_key=embedding_api_key,
    model= embedding_model,
    azure_endpoint=embedding_azure_endpoint,
    api_version=embedding_api_version
)

#Text Splitter Initialization
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=100, 
    length_function=len, 
    is_separator_regex=False
    )

#AZURECOSMOSDB Initialization
db_username = os.getenv('AZURECOSMOSDB_USERNAME')
db_password = os.getenv('AZURECOSMOSDB_PASSWORD')
username = quote_plus(db_username)
password = quote_plus(db_password)

uri = f"mongodb+srv://{username}:{password}@dxcragcosmodbcluster1.mongocluster.cosmos.azure.com/?tls=true&authMechanism=SCRAM-SHA-256&retrywrites=false&maxIdleTimeMS=120000"

client = MongoClient(uri)


DB_NAME= "embeddings"
db = client[DB_NAME]

COLLECTION_NAME = "embeddings_collection"
collection = client[DB_NAME][COLLECTION_NAME]
embedding_key = "contentVector"

INDEX_NAME = "VectorSearchIndex"

NAMESPACE = DB_NAME + "." + COLLECTION_NAME

if COLLECTION_NAME not in db.list_collection_names():
    # Creates a unsharded collection that uses the DBs shared throughput
    db.create_collection(COLLECTION_NAME)
    print("Created collection '{}'.\n".format(COLLECTION_NAME))
else:
    print("Using collection: '{}'.\n".format(COLLECTION_NAME))



################################################################################### PROMPTS ####################################################################################################
raw_prompt = PromptTemplate.from_template("""
    You are a technical assistant who is good at searching and analyzing documents. If you do not have an answer from the provided information then say so. 
    {input}
    Context: {context}
    Answer:                                     
""")


prompt = """
You are a power point presentation specialist. You are asked to create
the content for a presentation about {topic}.
You have been given the following information to create a presentation:
---
{information}.
---
Structure the information in a way that it can be put in a power point
presentation. Each slide should have a title and content, with the content
being a summary of the information provided. Each slide should have one or
more sentences that capture the key points of the information.
Return the structured information as a JSON as follows.
Your answer should only contain the JSON - no markdown formatting.
"""
prompt_template = PromptTemplate.from_template(prompt)
prompt_examples = """
Example:
{"slides": [
{"title": "Slide 1", "content": "Content for slide 1"},
{"title": "Slide 2", "content": "Content for slide 2"},
{"title": "Slide 3", "content": "Content for slide 3"},
]}
"""

powerpoint_prompt = """
You are a PowerPoint presentation specialist. You'll get a list of slides, each
slide containing a title and content. You need to create a detailed and insightful
PowerPoint presentation based on the provided slides.
Produce the PowerPoint slides in themes relevant to the topic to make them more interesting and captivating.
But there is a catch: Instead of creating the presentation, provide python code
that generates the PowerPoint presentation based on the provided slides.
Please use the package python-pptx to create the PowerPoint presentation.
The presentation should be visually appealing and professionally designed.
If the slides' content contains more than one information, make bullet points.
Make sure to initialize the Presentation object with `prs = Presentation()`.
Make sure to mport the `Pt` function from `pptx.util` for defining font sizes if it is used.
Make sure to Save the presentation as 'presentation.pptx' inside the path 'static/powerpoint/'.
Your answer must and should only contain the correct Python code, no explanatory text.
Slides:
"""


def ppt_code_checker(ppt_code):
    prompt = """
    Check that the following code is in the correct format by avoiding NameError: name 'prs' is not defined, Syntax Errors, and using the most updated version of python-pptx so that no errors occur during execution.  
    Ensure that `prs = Presentation()` is initialized if not already done.
    Ensure the presentation is saved as 'presentation.pptx' inside the path 'static/powerpoint/'. 
    Ensure to import the `Pt` function from `pptx.util` for defining font sizes if it is used.
    If the code is not in the correct format, correct it and generate the corrected code that can be executed. Any relevant libaries should be imported.
    Your answer must and should only contain the correct Python code, no explanatory text.
    Code:
    {code}
    """
    llm_prompt = (
        prompt.format(code=ppt_code)
    )
    correct_code = llm.invoke(llm_prompt)
    
    return correct_code.content


def generate_test_information(query):
    prompt = """
    Generate a detailed and informative passage with this response: {topic}. 
    The passage should provide an all-encompassing view of the topic, including the topic's importance, benefits, and any relevant advice or recommendations. 
    The language should be clear, concise, and suitable for a general audience. """
    
    llm_prompt = (
    prompt.format(topic=query)
    )

    test_information = llm.invoke(llm_prompt)
    
    return(test_information.content)

# def generate_powerpoint(query):
    
#     # Generating the slides content
    
#      # Check if a pptx file already exists
#     if os.path.exists('static/powerpoint/presentation.pptx'):
#         # Delete the file
#         os.remove('static/powerpoint/presentation.pptx')
    
#     content_prompt = (
#         prompt_template.format(topic=query, information=generate_test_information(query))
#         + prompt_examples
#     )
    
#     slides_response = llm.invoke(content_prompt)
#     slides = slides_response.content # Assuming response is a dictionary with 'content' key

#     # Generating the PowerPoint creation code
#     presentation_code_response = llm.invoke(powerpoint_prompt + slides)
#     presentation_code = presentation_code_response.content
    
#     correct_python_code = ppt_code_checker(presentation_code)
    
#     # Extracting the python code from the response
#     match = re.findall(r"python\n(.*?)\n```", correct_python_code, re.DOTALL)
#     python_code = match[0]

#     # Executing the extracted Python code
#     exec(python_code)


from pptx import Presentation
from pptx.util import Inches, Pt
import json
def generate_powerpoint(query):
    # Generate slide content using the LLM
    content_prompt = (
        prompt_template.format(topic=query, information=generate_test_information(query))
        + prompt_examples
    )
    
    slides_response = llm.invoke(content_prompt)
    slides_content_str = slides_response.content  # Get the raw string content
    
    try:
        # Parse the string content as JSON
        slides_content = json.loads(slides_content_str)
    except json.JSONDecodeError as e:
        print(f"Failed to parse slides content as JSON: {e}")
        return

    # Initialize the PowerPoint presentation
    prs = Presentation()
    
    # Iterate over the slides content and add them to the presentation
    for slide in slides_content['slides']:
        slide_layout = prs.slide_layouts[1]  # Choosing a layout with title and content
        slide_obj = prs.slides.add_slide(slide_layout)
        title = slide_obj.shapes.title
        content = slide_obj.shapes.placeholders[1].text_frame
        
        title.text = slide['title']
        
        # Clear any existing paragraphs in the content placeholder
        content.clear()

        # Add content as bullet points if necessary
        if isinstance(slide['content'], list):
            for item in slide['content']:
                p = content.add_paragraph()
                p.text = item
                p.font.size = Pt(18)
        else:
            p = content.add_paragraph()
            p.text = slide['content']
            p.font.size = Pt(18)
    
    # Check if a pptx file already exists
    if os.path.exists('static/powerpoint/presentation.pptx'):
        os.remove('static/powerpoint/presentation.pptx')
    
    # Save the presentation
    prs.save('static/powerpoint/presentation.pptx')
################################################################################################################################################################################################


@app.route("/upload_documents", methods=["POST"])
@cross_origin()
def uploadDocuments():
    files = request.files.getlist('file')  # Get list of files
    responses = []

    for file in files:
        file_extension = os.path.splitext(file.filename)[1].lower()
        filename = file.filename  # Initialize filename variable
        
        if file_extension == '.pdf':
            pdf_filename = file.filename
            # SAVE pdf TO static/pdf FOLDER
            pdf_path = os.path.join('static/pdf/uploads', pdf_filename)
            os.makedirs(os.path.dirname(pdf_path), exist_ok=True)
            file.save(pdf_path)
            
            path = 'static/pdf/uploads/' + pdf_filename

            loader = PyMuPDFLoader(path)
            docs = loader.load_and_split()
            print(f"docs len = {len(docs)}")

            chunks = text_splitter.split_documents(docs)
            print(f"chunks len = {len(chunks)}")
        
        elif file_extension == '.txt':
            txt_filename = file.filename
            # SAVE txt TO static/txt FOLDER
            txt_path = os.path.join('static/txt/uploads', txt_filename)
            os.makedirs(os.path.dirname(txt_path), exist_ok=True)
            file.save(txt_path)
            
            path = 'static/txt/uploads/' + txt_filename

            loader = TextLoader(path)
            docs = loader.load()
            print(f"docs len = {len(docs)}")

            chunks = text_splitter.split_documents(docs)
            print(f"chunks len = {len(chunks)}")

        vectorstore = AzureCosmosDBVectorSearch.from_documents(
                    chunks,
                    embedding_model,
                    collection=collection,
                    index_name = INDEX_NAME,
                    embedding_key=embedding_key,
                    )

        response = {"filename": filename, "doc_len": len(docs), "chunk_len": len(chunks)}
        responses.append(response)
        os.remove(path)
        
    num_lists = 100
    dimensions = 1536
    similarity_algorithm = CosmosDBSimilarityType.COS
    kind = CosmosDBVectorSearchType.VECTOR_IVF
    m = 16
    ef_construction = 64

    vectorstore.create_index(
        num_lists, dimensions, similarity_algorithm, kind, m, ef_construction
    )


    return jsonify({"status": "Successfully Uploaded", "files": responses})



@app.route("/generate_ppt", methods=["POST"])
@cross_origin()
def generatePowerpoint():
    data = request.get_json()
    query = data.get("query")
    
    print(f"Query: {query}")
    
    print("Loading Vector Store")

    vectorstore = AzureCosmosDBVectorSearch.from_connection_string(
    connection_string=uri,
    namespace=NAMESPACE,
    embedding=embedding_model,
    index_name = INDEX_NAME,
    embedding_key=embedding_key
    )
    
    if vectorstore.index_exists() == False:
        generate_powerpoint(query)
    
    else:

        retriever = vectorstore.as_retriever(
        search_kwargs={"k": 10}
        )
        
        # retriever_prompt = ChatPromptTemplate.from_messages(
        #     [
        #     MessagesPlaceholder(variable_name = "chat_history"),
        #     ("human","{input}"),
        #     ("human","Given the above conversation, generate a search query to lookup relevant documents in order to get information relevant to the conversation",),
        #     ]
        #                                                     )
        
        # history_aware_retriever = create_history_aware_retriever(
        #     llm = llm,
        #     retriever = retriever,
        #     prompt = retriever_prompt
        #     )
        
        document_chain = create_stuff_documents_chain(llm, raw_prompt)
        
        chain = create_retrieval_chain(retriever, document_chain)
        
        result = chain.invoke({"input" : query})
        
        response_answer = {"answer" : result["answer"]}
        
        generate_powerpoint(response_answer)
        
    return send_file(path_or_file='static/powerpoint/presentation.pptx', as_attachment=True, mimetype= "application/vnd.openxmlformats-officedocument.presentationml.presentation" ,download_name='presentation.pptx')
    
    
    


    
@app.route("/reset_ppt", methods=["POST"])
@cross_origin()
def deletePowerPoint():
    collection.drop()
    collection.drop_indexes()
    
    file_path = 'static/powerpoint/presentation.pptx'
    if os.path.exists(file_path):
        os.remove(file_path)

    return jsonify({"status": "Database Cleared, Powerpoint Deleted"})

    
if __name__ == '__main__':
    app.run(debug=True)